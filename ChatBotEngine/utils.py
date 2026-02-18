######## TOUTES LES FONCTIONS POUR ENVOYER LES INSTRUCTIONS A l'IA ET POUR INTERPRETER SES REPONSES


############ STRUCTURE DE REPONSE ATTENDUE

def _build_agent_directive(agent_content: str) -> str:
    base = r"""
{{AGENT_BLOCK}}

Tu vas recevoir des instructions qu'un utilisateur aura envoyé à une UX de Chat, appelée MASSY CHAT. La réponse que tu feras ne sera pas directement envoyée à l'utilisateur mais elle sera d'abord interprêtée par MASSY CHAT qui fera ensuite la réponse à l'utilisateur.
Tu peux, toi, donner des instructions à MASSY CHAT pour générer des fichiers si tu penses que cela répondra à l'utilisateur. L'UX sait générer des fichiers PDF, DOCX, XLSX, PPTX et TXT (via des outils appropriés ReportLab, OpenPyxl, python-docx, python-pptx...) mais uniquement ces formats.
Tu donneras tes instructions à l'UX via la structure JSON  "files" qui est décrite plus loin. Ne génère d'image que SI et UNIQUEMENT SI l'utilisateur te demande explicitement de le faire.


Massy CHAT ne sait traiter qu'une seule réponse. Tu ne dois donc fournir qu'un seul et unique JSON de réponse. Un seul JSON. Jamais deux. Jamais rien avant ce JSON, Jamais rien après. Aucun commentaire, aucune explication, aucune phrase d'introduction ou quoi que ce soit en dehors de ce qui est permis dans ce JSON.
Si tu ne respecte pas cette structure JSON, MASSY CHAT ne pourra pas interprêter ta réponse, tu feras bugger MASSY CHAT et l'utilisateur ne recevra aucune réponse.

Quelle sur soit réponse, elle doit donc OBLIGATOIREMENT être intégrée dans un et un seul bloc UNIQUE **JSON strictement valide** respectant STRICTEMENT la spécification suivante qui te permet de transmettre à la fois un texte explicatif et le contenu des éventuels fichiers :
{
  "main_markdown": "Texte en Markdown lisible",
  "files": [
    {
      "filename": "nom.ext",
      "type": "txt|xlsx|docx|pdf|pptx",
      "description": "Brève description",
      "content": <contenu dépendant du type>
    }
  ]
}

Aucun autre champ n’est autorisé à ce premier niveau. Aucun texte en dehors de cette structure JSON, jamais de ```json, jamais de texte avant ou après quel qu'il soit.

# Règles pour `main_markdown`

- doit être du Markdown pur (titres, listes, texte…)
- ne doit JAMAIS contenir de JSON
- ne doit PAS dupliquer le contenu des fichiers
- ne doit pas inclure d’explications sur la structure JSON

# Règles pour la structure "files" de ce JSON : il peut y avoir 0, 1 ou plusieurs structures. Ces structures sont destinées à stocker le contenu des éventuels fichiers qui feront partie de ta réponse et qui seront pris en charge par l'UX de cet agent. Chacune de ces structures NE DOIT contenir QUE :
- filename
- type
- description
- content

Rien d'autre.

---

# Spécifications du champ `content` par type de fichier

## 1. Type `txt`
- `content` est une chaîne de texte brut
- retours à la ligne `\n` autorisés
- jamais de markdown, jamais de json

---

## 2. Type `xlsx` — Fichier Excel qui sera ensuite généré via **openpyxl** par l'UX :

Le champ `content` doit être lui-même un **objet JSON strictement structuré**, compatible avec openpyxl :

{
  "sheets": [
    {
      "name": "Nom",
      "rows": [
        [
          "Texte",
          123,
          { "value": "Texte", "bold": true, "italic": true, "font_size": 14, "align": "center" }
        ],
        ...
      ],
      "column_widths": [ 25, 40, 15 ]
    }
  ]
}

### Formatage autorisé (ET UNIQUEMENT celui-ci)
Pour les objets cellule :

- "value": texte ou nombre  (obligatoire)
- "bold": true/false
- "italic": true/false
- "font_size": nombre (ex. 12, 14, 18)
- "align": "left" | "center" | "right"

### ❌ Restrictions strictes (IMPORTANT)
Parce que le fichier est généré via `openpyxl`, tu N’AS PAS LE DROIT d'utiliser :

- couleurs (ni texte ni fond)
- bordures
- fusion de cellules
- formules
- styles complexes
- mises en forme conditionnelles
- formats numériques avancés
- images
- hyperliens

Tu dois rester dans un format Excel simple, purement textuel.

---

## 3. Type `docx` — Document Word qui sera ensuite généré via **python-docx** par l'UX :

Le champ `content` doit être lui-même un **objet JSON strictement structuré**, compatible avec python-docx :

Structure obligatoire :

{
  "title": {
    "text": "Titre",
    "bold": true,
    "align": "center",
    "size": 20
  },
  "paragraphs": [
    {
      "text": "Texte du paragraphe",
      "italic": true,
      "align": "center",
      "size": 12
    }
  ]
}

### ❌ INTERDITS ABSOLUS
- JAMAIS de RTF (`{\rtf1 ...}`)
- JAMAIS de HTML
- JAMAIS de Markdown dans `content`
- JAMAIS de styles Word avancés non supportés
- JAMAIS d’images, tableaux, listes ou objets riches

---

#### 4. Type `pdf`

Le champ `content` doit être un objet JSON strict, destiné à être converti en PDF via la bibliothèque Python **ReportLab**.

La structure doit être :

{
  "title": "Titre du document (optionnel)",
  "paragraphs": [
    {
      "text": "Paragraphe de texte",
      "bold": true,
      "italic": true,
      "size": 14,
      "align": "center"
    }
  ]
}

Contraintes :
- Le PDF est généré avec ReportLab dans cette application.
- Only texte simple : pas d'images, pas de tableaux.
- Pas de couleurs, pas de bordures.
- Police : Helvetica uniquement.
- Alignements : left / center / right.
- Tout doit tenir dans ce JSON (pas de Markdown ou HTML dans les champs texte).

---

## 5. Type `pptx` — Présentation PowerPoint générée via python-pptx

Le champ `content` doit être un objet JSON structuré ainsi :

{
  "slides": [
    {
      "title": "Titre de la diapositive",
      "paragraphs": [
        { "text": "Ligne 1", "bold": true, "italic": false, "size": 24, "align": "center" }
      ]
    }
  ]
}

Contraintes :
- Pas d’images
- Pas de tableaux
- Pas de mise en page avancée
- Uniquement du texte simple
- Align : left / center / right
- Polices : Calibri uniquement

---

Voilà pour la STRUCTURE de ta réponse, à respecter absolument. Le CONTENU de ta réponse sera généré selon les instructions utilisateurs qui suivent : 

"""
    return base.replace("{{AGENT_BLOCK}}", agent_content.strip())


#################   EXTRACTION DU TEXTE DE REPONSE PRINCIPALE



def extract_markdown_from_agent_response(raw_response: str) -> str:
    if not isinstance(raw_response, str):
        return ""

    match = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", raw_response, flags=re.IGNORECASE)
    candidate = match.group(1).strip() if match else raw_response.strip()

    try:
        data = json.loads(candidate)
        if isinstance(data, dict) and "main_markdown" in data:
            return data["main_markdown"].strip()
    except Exception:
        try:
            inner_match = re.search(r'"main_markdown"\s*:\s*"([^"]+)"', candidate)
            if inner_match:
                return inner_match.group(1).strip()  # ✅ plus d'encodage/décodage
        except Exception:
            pass

    return candidate.strip()


##############   GENERATION DU TXT DE DEBUG DECRIVANT LA STRUCTURE DE LA REPONSE :

def describe_agent_json_structure(raw_text: str) -> str:
    """
    Décrit de manière lisible la structure JSON renvoyée par l'IA :
    - Markdown seul
    - Avec fichiers (txt, xlsx, docx, pdf)
    - Type(s) détectés
    Ne se base que sur une analyse textuelle (pas un parsing JSON).
    """
    if not raw_text or not raw_text.strip():
        return "⚪ Structure IA : vide"

    text = raw_text.strip()

    # 🌟 Première étape : repérer la présence des champs clés
    has_main = '"main_markdown"' in text
    has_files = '"files"' in text

    if has_main and has_files:
        # Vérification robuste des "files": []
        normalized = re.sub(r"\s+", "", text)
        if '"files":[]' in normalized:
            return "🟢 Réponse Markdown seule"

        # 🌟 Extraction des types déclarés
        file_types = re.findall(r'"type"\s*:\s*"([^"]+)"', text, flags=re.IGNORECASE)
        if not file_types:
            return "🟢 Markdown + fichiers (type non détecté)"

        labels = []
        for t in file_types:
            t_low = t.lower()
            if t_low == "txt":
                labels.append("fichier texte")
            elif t_low == "xlsx":
                labels.append("fichier Excel")
            elif t_low == "docx":
                labels.append("fichier Word")
            elif t_low == "pdf":
                labels.append("fichier PDF")
            elif t_low == "pptx":
                labels.append("fichier Powerpoint")
            else:
                labels.append(f"fichier {t}")

        # Évite les doublons dans l'affichage
        labels = list(dict.fromkeys(labels))
        return "🟢 Markdown + " + " + ".join(labels)

    # Cas : main_markdown présent mais pas files
    if has_main:
        return "🟢 Réponse Markdown seule"

    # Cas : absence totale de JSON structuré
    return "🟢 Réponse texte libre / Markdown direct"



####################   TRAITEMENT DES FICHIERS REçUS

import json, re
from pathlib import Path
from django.conf import settings
from openpyxl import Workbook


# FONCTION UTILITAIRE POUR TENIR COMPTE DES VARIATIONS DANS LE FORMAT DE LA REPONSE

import json

def _extract_json_candidate(raw: str) -> str:
    """
    Extrait le premier JSON valide présent dans la chaîne `raw`,
    même si plusieurs JSON complets apparaissent dans la réponse.
    
    Utilise json.JSONDecoder().raw_decode mais élimine les faux positifs :
    - cas où plusieurs JSON sont collés (Mistral multiplie les outputs)
    - cas où du texte suit immédiatement et rendrait le JSON invalide.
    """
    if not isinstance(raw, str):
        return ""

    s = raw.strip()
    decoder = json.JSONDecoder()

    for i in range(len(s)):
        try:
            obj, end = decoder.raw_decode(s[i:])

            # --- Vérification cruciale ---
            tail = s[i+end:].lstrip()

            # Si tail commence par '{' → c’est un DEUXIÈME JSON → on ignore celui-là
            if tail.startswith("{"):
                continue  # on ne retourne pas ce JSON, on poursuit la recherche

            # Si tail commence par un caractère alphanum (ou "text" etc.)
            # cela signifie que le JSON a été concaténé brutalement à du texte → on ignore.
            if tail and tail[0].isalnum():
                continue

            # --- JSON valide et isolé : on le retourne tel quel ---
            return s[i:i+end]

        except json.JSONDecodeError:
            continue

    return ""




# TRAITEMENT SEQUENTIEL DE TOUS LES FICHIERS

def extract_files_from_agent_response(raw_response: str, request=None) -> list[dict]:
    """
    Extrait les fichiers depuis le premier JSON valide présent dans la réponse Mistral.
    Fonction robuste :
    - ignore les JSON concaténés (rely on _extract_json_candidate)
    - ne casse jamais si le JSON est incomplet
    - gère txt / xlsx / docx
    """
    # 1️⃣ Extraire le JSON isolé
    candidate = _extract_json_candidate(raw_response)
    if not candidate:
        return []  # aucun JSON exploitable

    # 2️⃣ Charger le JSON
    try:
        data = json.loads(candidate)
    except Exception:
        return []

    if not isinstance(data, dict):
        return []

    files = data.get("files")
    if not isinstance(files, list) or not files:
        return []  # aucun fichier

    # 3️⃣ Préparer le dossier d’export
    export_dir = Path(settings.MEDIA_ROOT) / "chat_exports"
    export_dir.mkdir(parents=True, exist_ok=True)

    out = []

    # 4️⃣ Traiter chaque fichier
    for f in files:
        if not isinstance(f, dict):
            continue

        filename    = (f.get("filename") or "fichier.txt").strip()
        ftype       = (f.get("type") or "txt").strip().lower()
        description = f.get("description") or ""
        content     = f.get("content")

        try:
            # --- TXT ---
            if ftype == "txt" and isinstance(content, str):
                path = _generate_txt_file(filename, content, export_dir)

            # --- XLSX ---
            elif ftype == "xlsx" and isinstance(content, dict):
                path = _generate_xlsx_file(filename, content, export_dir)

            # --- DOCX ---
            elif ftype == "docx" and isinstance(content, (str, dict)):
                path = _generate_docx_file(filename, content, export_dir)
            
            # --- PDF ---
            elif ftype == "pdf" and isinstance(content, dict):
                path = _generate_pdf_file(filename, content, export_dir)
                
            # --- PPTX ---
            elif ftype == "pptx" and isinstance(content, dict):
                path = _generate_pptx_file(filename, content, export_dir)

            # --- Type non géré ---
            else:
                continue

            # Construction de l’URL publique
            url = (request.build_absolute_uri(
                settings.MEDIA_URL + f"chat_exports/{path.name}"
            ) if request else None)

            out.append({
                "filename": filename,
                "description": description,
                "url": url,
                "type": ftype,
            })

        except Exception:
            # On ignore ce fichier mais on n’arrête pas le lot
            continue

    return out


# GENERATION TXT

def _generate_txt_file(filename: str, content: str, export_dir: Path) -> Path:
    safe_name = re.sub(r"[^a-zA-Z0-9_\-.]", "_", filename)
    p = export_dir / safe_name
    p.write_text(content or "", encoding="utf-8")
    return p

from openpyxl import Workbook
from openpyxl.styles import Font, Alignment
from pathlib import Path
import re

def _generate_xlsx_file(filename: str, content: dict, export_dir: Path) -> Path:
    """
    Génère un fichier XLSX à partir d'une structure JSON enrichie :
    {
      "sheets": [
        {
          "name": "Poème",
          "rows": [
            [
              { "value": "...", "bold": true, "font_size": 16, "font": "Arial", "align": "center" }
            ],
            [],
            [
              { "value": "...", "italic": true }
            ]
          ],
          "column_widths": [60]
        }
      ]
    }
    """

    safe_name = re.sub(r"[^a-zA-Z0-9_\-.]", "_", filename)
    path = export_dir / safe_name

    wb = Workbook()
    wb.remove(wb.active)  # on supprime la feuille par défaut

    for sheet in content.get("sheets", []):
        name = sheet.get("name", "Feuille1")
        rows = sheet.get("rows", [])
        widths = sheet.get("column_widths", [])

        ws = wb.create_sheet(title=name)

        # -- Largeurs de colonnes --
        for idx, w in enumerate(widths, start=1):
            try:
                ws.column_dimensions[chr(64 + idx)].width = w
            except Exception:
                pass

        # -- Contenu cellule par cellule --
        for row in rows:
            # ligne vide
            if row == []:
                ws.append([""])
                continue

            excel_row = []
            style_specs = []

            for cell in row:
                if isinstance(cell, dict) and "value" in cell:
                    value = cell.get("value", "")
                    excel_row.append(value)
                    style_specs.append(cell)
                else:
                    # cellule simple ("texte brut")
                    excel_row.append(cell)
                    style_specs.append(None)

            # append row → openpyxl crée les cellules
            ws.append(excel_row)

            # récupère l'index de la dernière ligne écrite
            row_idx = ws.max_row

            # appliquer formatages
            for col_idx, spec in enumerate(style_specs, start=1):
                if not isinstance(spec, dict):
                    continue

                cell = ws.cell(row=row_idx, column=col_idx)

                # Font
                font_kwargs = {}
                if spec.get("bold"):
                    font_kwargs["bold"] = True
                if spec.get("italic"):
                    font_kwargs["italic"] = True
                if "font_size" in spec:
                    try:
                        font_kwargs["size"] = float(spec["font_size"])
                    except:
                        pass
                if "font" in spec:
                    font_kwargs["name"] = spec["font"]

                cell.font = Font(**font_kwargs)

                # Alignement
                align = spec.get("align", "").lower()
                if align in ("center", "right", "left"):
                    cell.alignment = Alignment(horizontal=align)

    wb.save(path)
    return path


#GENERATION DOC

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt

def _generate_docx_file(filename: str, content, export_dir: Path) -> Path:
    """
    Génère un fichier DOCX à partir :
    - soit d’une chaîne simple (texte brut)
    - soit d’un dict structuré contenant {title:{...}, paragraphs:[...]}
    """

    safe_name = re.sub(r"[^a-zA-Z0-9_\-.]", "_", filename)
    p = export_dir / safe_name

    doc = Document()

    # ---------- CAS 1 : texte brut ----------
    if isinstance(content, str):
        for line in content.split("\n"):
            doc.add_paragraph(line)
        doc.save(p)
        return p

    # ---------- CAS 2 : contenu structuré ----------
    if isinstance(content, dict):

        # Helper: normalisation align
        def parse_align(v):
            if not isinstance(v, str):
                return WD_ALIGN_PARAGRAPH.LEFT
            v = v.lower()
            return {
                "center": WD_ALIGN_PARAGRAPH.CENTER,
                "right": WD_ALIGN_PARAGRAPH.RIGHT,
                "justify": WD_ALIGN_PARAGRAPH.JUSTIFY
            }.get(v, WD_ALIGN_PARAGRAPH.LEFT)

        # Helper: ajout d’un paragraphe formaté
        def add_paragraph(block):
            text = block.get("text", "")
            para = doc.add_paragraph()
            run = para.add_run(text)

            # Bold / Italic
            if block.get("bold"):
                run.bold = True
            if block.get("italic"):
                run.italic = True

            # Taille
            if isinstance(block.get("size"), (int, float)):
                run.font.size = Pt(block["size"])

            # Police
            if isinstance(block.get("font"), str):
                run.font.name = block["font"]

            # Alignement
            para.alignment = parse_align(block.get("align"))

            # Interligne
            if block.get("line_spacing"):
                try:
                    para.paragraph_format.line_spacing = float(block["line_spacing"])
                except:
                    pass

        # ----- Titre -----
        title = content.get("title")
        if isinstance(title, dict):
            add_paragraph(title)

        # ----- Paragraphes -----
        for block in content.get("paragraphs", []):
            if isinstance(block, dict):
                add_paragraph(block)

        doc.save(p)
        return p

    # Aucun format reconnu
    raise ValueError("Format DOCX non reconnu")

# GENERATION PDF

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from pathlib import Path
import re

def _generate_pdf_file(filename: str, content: dict, export_dir: Path) -> Path:
    """
    Génère un PDF simple avec ReportLab :
    - titre optionnel
    - paragraphes simples
    - align : left/center/right
    - bold/italic/size
    """

    safe_name = re.sub(r"[^a-zA-Z0-9_\-.]", "_", filename)
    p = export_dir / safe_name

    doc = SimpleDocTemplate(str(p), pagesize=A4, leftMargin=2*cm, rightMargin=2*cm)
    styles = getSampleStyleSheet()

    flow = []

    # Mapping alignements
    ALIGN = {
        "left": TA_LEFT,
        "center": TA_CENTER,
        "right": TA_RIGHT
    }

    # --- Titre ---
    title = content.get("title")
    if isinstance(title, str):
        style = ParagraphStyle(
            "title_style",
            parent=styles["Heading1"],
            alignment=TA_CENTER,
            fontName="Helvetica-Bold"
        )
        flow.append(Paragraph(title, style))
        flow.append(Spacer(1, 12))

    # --- Paragraphes ---
    for pdef in content.get("paragraphs", []):
        if not isinstance(pdef, dict):
            continue

        text = pdef.get("text", "")
        size = pdef.get("size", 12)
        bold = pdef.get("bold", False)
        italic = pdef.get("italic", False)
        align = pdef.get("align", "left")

        font = "Helvetica"
        if bold and italic:
            font = "Helvetica-BoldOblique"
        elif bold:
            font = "Helvetica-Bold"
        elif italic:
            font = "Helvetica-Oblique"

        style = ParagraphStyle(
            "p_style",
            parent=styles["BodyText"],
            fontName=font,
            fontSize=size,
            alignment=ALIGN.get(align, TA_LEFT)
        )

        flow.append(Paragraph(text, style))
        flow.append(Spacer(1, 12))

    doc.build(flow)
    return p

# GENERATION PPTX

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

def _generate_pptx_file(filename: str, content: dict, export_dir: Path) -> Path:
    """
    Génère un fichier PPTX simple :
    {
      "slides": [
        {
          "title": "Titre",
          "paragraphs": [
            { "text": "...", "bold": true, "italic": false, "size": 24, "align": "center" }
          ]
        }
      ]
    }
    """
    safe_name = re.sub(r"[^a-zA-Z0-9_\-.]", "_", filename)
    path = export_dir / safe_name

    prs = Presentation()

    for slide_def in content.get("slides", []):
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)

        # ---- Titre ----
        if "title" in slide_def and slide.shapes.title:
            slide.shapes.title.text = slide_def["title"]

        # ---- Trouver la zone de contenu ----
        body = None
        for ph in slide.placeholders:
            if ph.is_placeholder and ph.placeholder_format.type == 1:  # BODY
                body = ph.text_frame
                break

        if body is None:
            # fallback : textbox manuelle
            box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            body = box.text_frame

        body.text = ""  # nettoie le contenu

        # ---- Paragraphes ----
        for pdef in slide_def.get("paragraphs", []):
            text = pdef.get("text", "").strip()
            if not text:
                continue

            size = int(pdef.get("size", 20))
            bold = bool(pdef.get("bold", False))
            italic = bool(pdef.get("italic", False))
            align = pdef.get("align", "left").lower()

            para = body.add_paragraph()
            run = para.add_run()
            run.text = text

            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic

            para.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT
            }.get(align, PP_ALIGN.LEFT)

    prs.save(path)
    return path
